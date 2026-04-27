from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import urllib.request
import tempfile

mcp = FastMCP("ppt_server")

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Design constants
BG_COLOR        = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT_COLOR     = RGBColor(0x00, 0xB4, 0xD8)   # cyan
TITLE_COLOR      = RGBColor(0xFF, 0xFF, 0xFF)   # white
BODY_COLOR       = RGBColor(0xCA, 0xD3, 0xE0)   # light grey
HIGHLIGHT_COLOR  = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent


def resolve(filename: str) -> str:
    if os.path.isabs(filename):
        return filename
    return os.path.join(BASE_DIR, "..", filename)


def set_slide_background(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=BODY_COLOR,
                align=PP_ALIGN.LEFT, wrap=True):
    """Helper to add a styled textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def download_image(query: str) -> str:
    """Download an image from Pexels for the given search query."""
    try:
        import urllib.parse
        import urllib.request
        import json
        import tempfile
        import os
        import sys
        sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
        from config import Config

        encoded = urllib.parse.quote(query)
        url = f"https://api.pexels.com/v1/search?query={encoded}&per_page=1&orientation=landscape"

        req = urllib.request.Request(
            url,
            headers={
                "Authorization": Config.PEXELS_API_KEY,
                "User-Agent": "AutoPPT/1.0"
            }
        )

        with urllib.request.urlopen(req, timeout=10) as response:
            data = json.loads(response.read())

        photos = data.get("photos", [])
        if not photos:
            print(f"   No image found for query: {query}")
            return None

        img_url = photos[0]["src"]["large"]
        print(f"   Downloading image: {img_url}")

        img_req = urllib.request.Request(
            img_url,
            headers={"User-Agent": "AutoPPT/1.0"}
        )
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        with urllib.request.urlopen(img_req, timeout=15) as img_response:
            tmp.write(img_response.read())
        tmp.close()

        print(f"   Image saved to {tmp.name}")
        return tmp.name

    except Exception as e:
        print(f"  Image download failed: {e}")
        return None


@mcp.tool()
def plan_slides(topic: str, num_slides: int = 5) -> str:
    """
    Plan the slide structure before creating the presentation.
    Always call this first.

    Args:
        topic: Subject of the presentation
        num_slides: Number of slides (default 5)

    Returns:
        Confirmation to proceed
    """
    return (
        f"Plan acknowledged for '{topic}' with {num_slides} slides. "
        "Structure: Slide 1 = intro (no image), "
        "Slides 2-4 = content with images, "
        "Slide 5 = conclusion (no image). "
        "Now call create_presentation."
    )


@mcp.tool()
def create_presentation(filename: str) -> str:
    """
    Create a new blank PowerPoint presentation with dark theme.

    Args:
        filename: Output .pptx filename

    Returns:
        Success message
    """
    path = resolve(filename)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    prs.save(path)
    return f"Presentation '{filename}' created successfully."


@mcp.tool()
def add_intro_slide(filename: str, title: str, subtitle: str) -> str:
    """
    Add a styled intro/title slide with no image.
    Use this for the FIRST slide only.

    Args:
        filename: Target .pptx file
        title: Main presentation title
        subtitle: Short subtitle or description

    Returns:
        Success message
    """
    path = resolve(filename)
    if not os.path.exists(path):
        return f"Error: '{filename}' not found. Call create_presentation first."

    prs = Presentation(path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_background(slide, BG_COLOR)

    W = prs.slide_width
    H = prs.slide_height

    # Accent bar on left
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0), Inches(0),
        Inches(0.15), H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Title
    add_textbox(slide, title,
                left=Inches(0.5), top=Inches(2.2),
                width=Inches(12), height=Inches(1.5),
                font_size=44, bold=True,
                color=TITLE_COLOR, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, subtitle,
                left=Inches(0.5), top=Inches(3.8),
                width=Inches(12), height=Inches(1),
                font_size=22, bold=False,
                color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    prs.save(path)
    return f"Intro slide '{title}' added."


@mcp.tool()
def add_content_slide(filename: str, title: str,
                      bullet_points: list[str], image_query: str) -> str:
    """
    Add a content slide with title, bullet points on the left,
    and an auto-fetched image on the right.
    Use this for slides 2, 3, and 4.

    Args:
        filename: Target .pptx file
        title: Slide title
        bullet_points: List of 3-5 bullet point strings
        image_query: Search term for the slide image (e.g. 'star nebula space')

    Returns:
        Success message
    """
    path = resolve(filename)
    if not os.path.exists(path):
        return f"Error: '{filename}' not found."

    prs = Presentation(path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    W = prs.slide_width
    H = prs.slide_height

    # Accent top bar
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Title
    add_textbox(slide, title,
                left=Inches(0.4), top=Inches(0.15),
                width=Inches(7.5), height=Inches(0.8),
                font_size=28, bold=True,
                color=TITLE_COLOR)

    # Bullet points on left half
    bullet_text = "\n\n".join(f"• {bp}" for bp in bullet_points)
    add_textbox(slide, bullet_text,
                left=Inches(0.4), top=Inches(1.1),
                width=Inches(6.5), height=Inches(5.8),
                font_size=22, bold=False,
                color=BODY_COLOR)

    # Image on right half
    img_path = download_image(image_query)
    if img_path:
        try:
            slide.shapes.add_picture(
                img_path,
                left=Inches(7.2), top=Inches(0.9),
                width=Inches(5.7), height=Inches(5.8)
            )
            os.unlink(img_path)
        except Exception:
            pass  # if image fails, slide still looks fine

    prs.save(path)
    return f"Content slide '{title}' added with image."


@mcp.tool()
def add_conclusion_slide(filename: str, summary: str,
                          takeaways: list[str]) -> str:
    """
    Add a styled conclusion slide with no image.
    Use this for the LAST slide only.

    Args:
        filename: Target .pptx file
        summary: One sentence wrap-up of the presentation
        takeaways: List of 3 key takeaway strings

    Returns:
        Success message
    """
    path = resolve(filename)
    if not os.path.exists(path):
        return f"Error: '{filename}' not found."

    prs = Presentation(path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    W = prs.slide_width
    H = prs.slide_height

    # Bottom accent bar
    bar = slide.shapes.add_shape(
        1,
        Inches(0), H - Inches(0.08),
        W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Title
    add_textbox(slide, "Conclusion",
                left=Inches(0.5), top=Inches(0.3),
                width=Inches(12), height=Inches(0.9),
                font_size=36, bold=True,
                color=TITLE_COLOR, align=PP_ALIGN.CENTER)

    # Summary
    add_textbox(slide, summary,
                left=Inches(1), top=Inches(1.4),
                width=Inches(11), height=Inches(1),
                font_size=19, bold=False,
                color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    # Key takeaways
    takeaway_text = "\n".join(f"  {t}" for t in takeaways)
    add_textbox(slide, takeaway_text,
                left=Inches(2), top=Inches(2.6),
                width=Inches(9), height=Inches(3.5),
                font_size=18, bold=False,
                color=BODY_COLOR, align=PP_ALIGN.LEFT)

    prs.save(path)
    return "Conclusion slide added."


@mcp.tool()
def save_presentation(filename: str) -> str:
    """
    Confirm the presentation is saved on disk.
    Call this as the final step.

    Args:
        filename: The .pptx filename to confirm

    Returns:
        Success message with file path
    """
    path = resolve(filename)
    if not os.path.exists(path):
        return f"Error: '{filename}' not found."
    return f"Presentation '{filename}' saved successfully at {path}."


if __name__ == "__main__":
    mcp.run()