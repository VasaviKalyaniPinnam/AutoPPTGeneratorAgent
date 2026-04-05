import asyncio
import json
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage, SystemMessage, ToolMessage
from langchain_mcp_adapters.client import MultiServerMCPClient
from config import Config


def get_llm():
    return ChatGroq(
        model=Config.MODEL_ID,
        api_key=Config.GROQ_API_KEY,
        temperature=0,
    )


async def run_agent(user_prompt: str):
    # Handle vague prompts
    if len(user_prompt.strip()) < 5:
        user_prompt = "Create a general informative presentation"
        print(f"  Vague prompt detected. Using default: '{user_prompt}'")

    print(f"\n Starting AutoPPT Agent for topic: '{user_prompt}'\n")

    try:
        client = MultiServerMCPClient({
            "ppt_server": {
                "command": "python",
                "args": ["mcp_servers/ppt_server.py"],
                "transport": "stdio"
            },
            "search_server": {
                "command": "python",
                "args": ["mcp_servers/search_server.py"],
                "transport": "stdio"
            }
        })

        tools = await client.get_tools()
        print(f" Loaded {len(tools)} tools: {[t.name for t in tools]}\n")

        llm = get_llm()
        llm_with_tools = llm.bind_tools(tools, tool_choice="any")
        tool_map = {t.name: t for t in tools}

        messages = [
            SystemMessage(content=(
                "You are a PowerPoint presentation agent. You MUST follow these steps in order:\n\n"
                "STEP 1 - PLAN: Call plan_slides with the topic to outline the presentation.\n"
                "STEP 2 - SEARCH: Call search_web to get real information about the topic.\n"
                "STEP 3 - CREATE: Call create_presentation with filename='" + Config.OUTPUT_FILE + "'.\n"
                "STEP 4 - WRITE: Call add_slide 5 times, one per slide, using info from search results.\n"
                "STEP 5 - SAVE: Call save_presentation with filename='" + Config.OUTPUT_FILE + "'.\n\n"
                "Rules:\n"
                "- Never skip the planning step.\n"
                "- Always call tools. Never just describe what you would do.\n"
                "- If search fails, use your own knowledge — never crash.\n"
                "- Each slide must have a clear title and 3-5 bullet points.\n"
            )),
            HumanMessage(content=(
                f"Create a 5-slide PowerPoint presentation about: {user_prompt}\n\n"
                f"Output filename: {Config.OUTPUT_FILE}\n\n"
                "Begin with STEP 1: call plan_slides now."
            ))
        ]

        max_iterations = 30
        save_done = False

        for iteration in range(1, max_iterations + 1):
            print(f" Iteration {iteration}...")

            try:
                response = await llm_with_tools.ainvoke(messages)
            except Exception as llm_error:
                print(f" LLM error: {llm_error}")
                print(" Retrying with simplified message...")
                messages.append(HumanMessage(content="Continue from where you left off. Call the next tool."))
                continue

            messages.append(response)

            if not response.tool_calls:
                print("\n Agent finished.\n")
                print("Final message:", response.content)
                break

            for tool_call in response.tool_calls:
                tool_name = tool_call["name"]
                tool_args = tool_call["args"]
                tool_id   = tool_call["id"]

                print(f"   Calling: {tool_name}({json.dumps(tool_args)})")

                try:
                    if tool_name in tool_map:
                        result = await tool_map[tool_name].ainvoke(tool_args)
                    else:
                        result = f"Tool '{tool_name}' not found. Available tools: {list(tool_map.keys())}"
                except Exception as tool_error:
                    result = f"Tool error: {tool_error}. Continue with next step."

                print(f"    Result: {result}\n")

                messages.append(ToolMessage(
                    content=str(result),
                    tool_call_id=tool_id
                ))

                if tool_name == "save_presentation":
                    save_done = True
                    llm_with_tools = llm.bind_tools(tools)  # remove forced tool_choice

            if save_done:
                response = await llm_with_tools.ainvoke(messages)
                print("\n Done! Final message:", response.content)
                break

        else:
            print(" Max iterations reached. Attempting emergency save...")
            try:
                result = await tool_map["save_presentation"].ainvoke({"filename": Config.OUTPUT_FILE})
                print(f"  ✔  Emergency save result: {result}")
            except Exception as e:
                print(f" Emergency save failed: {e}")

    except Exception as e:
        print(f"\n Agent crashed: {e}")
        print("Creating fallback presentation...")
        try:
            from pptx import Presentation as PPTPresentation
            prs = PPTPresentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = user_prompt
            slide.placeholders[1].text_frame.paragraphs[0].text = "Content generation failed. Please try again."
            prs.save(Config.OUTPUT_FILE)
            print(f" Fallback presentation saved as '{Config.OUTPUT_FILE}'")
        except Exception as fallback_error:
            print(f" Fallback also failed: {fallback_error}")

    print(f"\n Presentation saved as '{Config.OUTPUT_FILE}'\n")